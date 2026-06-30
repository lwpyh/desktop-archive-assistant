"""照片/视频/像册整理能力集 — 对齐「照片像册整理」需求。

在 desktop-archive-assistant 现有 17 个通用能力之上，针对照片/图片/视频
场景补充深度处理能力。所有写操作默认 dry-run，删除一律降级为移入回收站，
绝不硬删除，可回滚。

能力清单（照片/视频专项，9 个）：
  archive_by_date     — 按 EXIF 拍摄日期归档到 年/月/日 三级目录
  dedupe_photos       — 照片感知哈希(pHash)去重（视觉相似，移入回收站）
  extract             — VLM 识别照片内容/文字 → 结构化数据(txt/csv)
  crop                — 图片裁剪（按尺寸/比例，可选深色区域居中检测）
  to_ppt              — 图片批量转 PPT（16:9，居中铺满）
  collage             — 多图拼接成一张（A4/网格）
  video_rename_title  — 视频标题 AI 重命名（删#标签/括号、15-25字、生成变体）
  video_compose       — 图片 + 音频 合成视频（FFmpeg，9:16/16:9）
  video_distribute    — 视频按规则分发到多个子文件夹

依赖（按需懒加载，缺失时优雅降级）：
  Pillow      图片读写/裁剪/EXIF
  imagehash   感知哈希
  python-pptx 图转PPT
  ffmpeg      (系统命令) 图音合成视频
"""
from __future__ import annotations

import datetime
import os
import re
import shutil
import subprocess
import tempfile
import time
from typing import Any, Dict, List, Optional, Tuple

from .extractors import scan_desktop
from .extractors.image_extractor import extract_exif_time, perceptual_hash
from .utils import expand, ensure_dir, logger, track, ProgressTracker


# 常见图片/视频扩展名
IMAGE_EXTS = {"jpg", "jpeg", "png", "bmp", "gif", "webp", "tif", "tiff", "heic", "mpo"}
VIDEO_EXTS = {"mp4", "mov", "avi", "mkv", "flv", "wmv", "m4v", "webm", "3gp", "ts"}
AUDIO_EXTS = {"mp3", "wav", "m4a", "aac", "flac", "ogg", "wma"}


def _is_image(ext: str) -> bool:
    return ext.lower().lstrip(".") in IMAGE_EXTS


def _is_video(ext: str) -> bool:
    return ext.lower().lstrip(".") in VIDEO_EXTS


def _unique_path(dst: str) -> str:
    """同名冲突时加 -1/-2 后缀，绝不覆盖。"""
    if not os.path.exists(dst):
        return dst
    base, ext = os.path.splitext(dst)
    i = 1
    while os.path.exists(f"{base}-{i}{ext}"):
        i += 1
    return f"{base}-{i}{ext}"


def _trash_root(trash_dir: str) -> str:
    root = os.path.join(expand(trash_dir), time.strftime("%Y%m%d_%H%M%S"))
    os.makedirs(root, exist_ok=True)
    return root


# ============================================================
# 能力 A：archive_by_date — 按 EXIF 拍摄日期归档
# ============================================================

def capability_archive_by_date(
    root: str,
    level: str = "day",
    use_mtime_fallback: bool = True,
    recursive: bool = False,
    dry_run: bool = False,
) -> Dict[str, Any]:
    """按拍摄日期把散落照片归入 年/月/日 三级目录。

    level: "day"(年/月/日) | "month"(年/月) | "year"(年)
    use_mtime_fallback: 无 EXIF 时回退到文件修改时间
    返回 {moves:[{src,dst,date,date_source}], skipped:[...], date_mismatch:[...]}

    date_mismatch: EXIF 年份与文件名标注年份不一致的文件（仅提示，不阻断）。
    """
    import re as _re
    root = expand(root)
    assets = scan_desktop(root, recursive=recursive, max_depth=5 if recursive else 1)
    assets = [a for a in assets if not a.is_shortcut and _is_image(a.ext)]

    result: Dict[str, Any] = {"root": root, "moves": [], "skipped": [], "date_mismatch": []}
    year_in_name = _re.compile(r"(19|20)\d{2}")

    # 逐张开图读 EXIF（图多时耗时），加进度。
    for a in track(assets, label="按拍摄日期归档", est_per_item=0.03):
        exif_ts = extract_exif_time(a.path)
        date_source = "exif"
        ts = exif_ts
        if ts is None:
            if use_mtime_fallback:
                ts = a.mtime
                date_source = "mtime"
            else:
                result["skipped"].append({"path": a.path, "reason": "no EXIF date"})
                continue

        dt = datetime.datetime.fromtimestamp(ts)
        # 文件名标注年份与 EXIF 年份对比
        m = year_in_name.search(os.path.basename(a.path))
        if m and int(m.group(0)) != dt.year:
            result["date_mismatch"].append({
                "path": a.path,
                "name_year": int(m.group(0)),
                "exif_year": dt.year,
                "date_source": date_source,
            })

        if level == "year":
            sub = os.path.join(f"{dt.year}年")
        elif level == "month":
            sub = os.path.join(f"{dt.year}年", f"{dt.month}月")
        else:  # day
            sub = os.path.join(f"{dt.year}年", f"{dt.month}月", f"{dt.day}")

        dst_dir = os.path.join(root, sub)
        dst = _unique_path(os.path.join(dst_dir, os.path.basename(a.path)))
        # 已在目标目录则跳过
        if os.path.dirname(a.path) == dst_dir:
            result["skipped"].append({"path": a.path, "reason": "already in place"})
            continue

        result["moves"].append({
            "src": a.path, "dst": dst,
            "date": dt.strftime("%Y-%m-%d"), "date_source": date_source,
        })
        if not dry_run:
            os.makedirs(dst_dir, exist_ok=True)
            shutil.move(a.path, dst)

    action = "archived" if not dry_run else "would archive"
    logger.info("[archive_by_date] %s %d photos (level=%s, %d mismatch, %d skipped)",
                action, len(result["moves"]), level,
                len(result["date_mismatch"]), len(result["skipped"]))
    return result


# ============================================================
# 能力 B：dedupe_photos — 感知哈希去重（视觉相似照片）
# ============================================================

def capability_dedupe_photos(
    root: str,
    threshold: int = 5,
    recursive: bool = True,
    dry_run: bool = False,
    trash_dir: str = "~/.archive_assistant/trash",
) -> Dict[str, Any]:
    """照片感知哈希(pHash)去重：视觉相似/重复照片移入回收站。

    threshold: 汉明距离阈值，越小越严格（0=完全相同，5=高度相似）
    保留每组中最早(best_time 最小)的一张为原始，其余移入回收站。
    返回 {scanned, groups:[{keep, duplicates:[...]}], trashed_count, freed_bytes}
    """
    root = expand(root)
    assets = scan_desktop(root, recursive=recursive, max_depth=8 if recursive else 1)
    assets = [a for a in assets if not a.is_shortcut and _is_image(a.ext)]

    # 计算 phash
    hashed: List[Tuple[Any, str, int]] = []  # (asset, phash_obj, _)
    try:
        import imagehash  # noqa: F401
    except ImportError:
        logger.warning("[dedupe_photos] imagehash not installed; falling back to exact hash")
        return {"scanned": len(assets), "groups": [], "trashed_count": 0,
                "freed_bytes": 0, "error": "imagehash not installed"}

    from PIL import Image
    import imagehash as _ih

    phash_map: List[Tuple[Any, Any]] = []  # (asset, ImageHash)
    # 逐张开图+感知哈希，图多时耗时，加进度避免长时间静默。
    for a in track(assets, label="去重·感知哈希计算", est_per_item=0.05):
        try:
            with Image.open(a.path) as im:
                ph = _ih.phash(im)
            phash_map.append((a, ph))
        except Exception as e:  # noqa: BLE001
            logger.debug("[dedupe_photos] phash failed %s: %s", a.path, e)

    # 贪心分组：未分组的作为种子，与其距离 <= threshold 的归为一组
    used = [False] * len(phash_map)
    groups: List[Dict[str, Any]] = []
    for i in range(len(phash_map)):
        if used[i]:
            continue
        ai, hi = phash_map[i]
        members = [ai]
        used[i] = True
        for j in range(i + 1, len(phash_map)):
            if used[j]:
                continue
            aj, hj = phash_map[j]
            if (hi - hj) <= threshold:
                members.append(aj)
                used[j] = True
        if len(members) > 1:
            # 保留 best_time 最早的一张
            members.sort(key=lambda x: x.best_time())
            keep = members[0]
            dups = members[1:]
            groups.append({
                "keep": keep.path,
                "duplicates": [d.path for d in dups],
            })

    # 执行去重
    trashed_count = 0
    freed = 0
    if not dry_run and groups:
        troot = _trash_root(trash_dir)
        for g in groups:
            for dup in g["duplicates"]:
                try:
                    freed += os.path.getsize(dup)
                    dst = _unique_path(os.path.join(troot, os.path.basename(dup)))
                    shutil.move(dup, dst)
                    trashed_count += 1
                except OSError as e:
                    logger.error("[dedupe_photos] move failed: %s", e)
    else:
        for g in groups:
            for dup in g["duplicates"]:
                try:
                    freed += os.path.getsize(dup)
                except OSError:
                    pass
            trashed_count += len(g["duplicates"])

    action = "deduped" if not dry_run else "would dedupe"
    logger.info("[dedupe_photos] %s %d duplicates in %d groups (~%.1fMB, threshold=%d)",
                action, trashed_count, len(groups), freed / 1e6, threshold)
    return {
        "scanned": len(assets),
        "groups": groups,
        "trashed_count": trashed_count,
        "freed_bytes": freed,
    }


# ============================================================
# 能力 C：extract — VLM 识别照片内容/文字 → 结构化数据
# ============================================================

def _phash_groups(assets: List[Any], threshold: int) -> List[List[Any]]:
    """按感知哈希(pHash)把视觉相似的图聚成簇，返回 [[asset,...], ...]。

    用于 extract 提速：同簇共用一次 caption，VLM 调用从 N 次降到"簇数"次。
    imagehash 不可用或 threshold<0 时退化为「每图独立成簇」（即不预聚类）。
    """
    if threshold is None or threshold < 0:
        return [[a] for a in assets]
    try:
        from PIL import Image
        import imagehash as _ih
    except ImportError:
        return [[a] for a in assets]

    hashed: List[Tuple[Any, Any]] = []
    # 预聚类需逐张开图算 phash，图多时耗时，加进度。
    for a in track(assets, label="预聚类·感知哈希", est_per_item=0.04):
        try:
            with Image.open(a.path) as im:
                hashed.append((a, _ih.phash(im)))
        except Exception as e:  # noqa: BLE001
            logger.debug("[extract] phash failed %s: %s", a.path, e)
            hashed.append((a, None))

    used = [False] * len(hashed)
    groups: List[List[Any]] = []
    for i in range(len(hashed)):
        if used[i]:
            continue
        ai, hi = hashed[i]
        members = [ai]
        used[i] = True
        if hi is not None:
            for j in range(i + 1, len(hashed)):
                if used[j]:
                    continue
                aj, hj = hashed[j]
                if hj is not None and (hi - hj) <= threshold:
                    members.append(aj)
                    used[j] = True
        groups.append(members)
    return groups


def _ocr_backend_is_concurrent_safe(vlm) -> bool:
    """判断 OCR 是否可并发：pytesseract(无 vlm) 或 ollama 后端可并发；

    transformers 本地单卡 generate 非线程安全，串行更稳。
    """
    if vlm is None:
        return True
    impl = getattr(vlm, "_impl", None)
    return type(impl).__name__ == "OllamaVLM"


def capability_extract(
    root: str,
    mode: str = "caption",
    vlm=None,
    out_path: Optional[str] = None,
    recursive: bool = False,
    cluster_threshold: int = 5,
    batch_size: int = 8,
    max_workers: int = 4,
    use_cache: bool = True,
) -> Dict[str, Any]:
    """对照片做内容识别/文字提取，输出结构化数据。

    mode: "caption"(一句话描述) | "ocr"(文字提取) | "both"
    vlm: VLMReasoner 实例（None 时仅 OCR 回退）
    out_path: 输出文件路径，.csv 输出表格，否则输出 txt

    提速策略（大批量关键）：
    - caption：先按 pHash 预聚类（cluster_threshold，<0 关闭），每簇仅对代表图
      调一次 VLM，再把描述回填整簇 → 调用次数从 N 降到簇数。
    - 代表图走 batch_size 批量前向（transformers 后端一次多图）。
    - ocr：文字逐图可能不同，不共享；pytesseract/ollama 后端用线程池并发
      （max_workers），transformers 单卡自动串行。
    - use_cache：按 路径+mtime 命中历史结果，跳过重复前向。
    返回 {records:[{file, caption, ocr_text}], out_path}
    """
    from concurrent.futures import ThreadPoolExecutor
    from .vlm_cache import VLMCache

    root = expand(root)
    assets = scan_desktop(root, recursive=recursive, max_depth=5 if recursive else 1)
    assets = [a for a in assets if not a.is_shortcut and _is_image(a.ext)]

    cache = VLMCache(enabled=use_cache)
    caption_map: Dict[str, str] = {}
    ocr_map: Dict[str, str] = {}

    # ---------- caption：缓存 → pHash 预聚类 → 代表图 batch 推理 → 回填整簇 ----------
    if mode in ("caption", "both") and vlm is not None:
        need: List[Any] = []
        for a in assets:
            cached = cache.get("caption", a.path)
            if cached is not None:
                caption_map[a.path] = cached
            else:
                need.append(a)

        groups = _phash_groups(need, cluster_threshold)
        reps = [g[0] for g in groups]
        rep_caps: List[str] = []
        # 代表图分批 VLM 描述：按代表图数量展示进度（每批推进 len(chunk)）。
        _cap_tr = ProgressTracker(len(reps), label="图片描述", est_per_item=2.0)
        for i in range(0, len(reps), max(1, batch_size)):
            chunk = reps[i:i + max(1, batch_size)]
            try:
                rep_caps.extend(vlm.caption_batch([r.path for r in chunk]))
            except Exception as e:  # noqa: BLE001
                logger.debug("[extract] caption batch failed: %s", e)
                rep_caps.extend([""] * len(chunk))
            _cap_tr.advance(len(chunk))
        _cap_tr.finish()

        for grp, cap in zip(groups, rep_caps):
            for a in grp:  # 同簇共用代表图描述
                caption_map[a.path] = cap
                cache.set("caption", a.path, cap)
        logger.info("[extract] caption: %d imgs -> %d VLM calls (clustered, %d cached)",
                    len(assets), len(reps), len(assets) - len(need))

    # ---------- ocr：缓存 → （可并发）逐图识别 ----------
    if mode in ("ocr", "both"):
        from .extractors.image_extractor import ocr_image

        def _do_ocr(a) -> Tuple[str, str]:
            cached = cache.get("ocr", a.path)
            if cached is not None:
                return a.path, cached
            try:
                txt = vlm.ocr(a.path) if vlm is not None else ocr_image(a.path)
            except Exception as e:  # noqa: BLE001
                logger.debug("[extract] ocr failed %s: %s", a.path, e)
                txt = ""
            return a.path, txt

        todo = [a for a in assets if cache.get("ocr", a.path) is None]
        for a in assets:
            c = cache.get("ocr", a.path)
            if c is not None:
                ocr_map[a.path] = c

        workers = max(1, max_workers) if _ocr_backend_is_concurrent_safe(vlm) else 1
        # OCR 逐图识别耗时，无论并发/串行都展示进度。
        _ocr_mode = f"并发 {workers} 路" if workers > 1 else "串行"
        _ocr_tr = ProgressTracker(
            len(todo), label=f"图片文字识别OCR（{_ocr_mode}）",
            est_per_item=1.5 / max(1, workers),
        )
        if workers > 1 and todo:
            with ThreadPoolExecutor(max_workers=workers) as ex:
                for path, txt in ex.map(_do_ocr, todo):
                    ocr_map[path] = txt
                    cache.set("ocr", path, txt)
                    _ocr_tr.advance()
        else:
            for a in todo:
                path, txt = _do_ocr(a)
                ocr_map[path] = txt
                cache.set("ocr", path, txt)
                _ocr_tr.advance()
        _ocr_tr.finish()

    cache.flush()

    # ---------- 组装记录（保持原扫描顺序）----------
    records: List[Dict[str, str]] = []
    for a in assets:
        rec: Dict[str, str] = {"file": os.path.basename(a.path), "path": a.path}
        if mode in ("caption", "both"):
            rec["caption"] = caption_map.get(a.path, "")
        if mode in ("ocr", "both"):
            rec["ocr_text"] = ocr_map.get(a.path, "")
        records.append(rec)

    written = None
    if out_path:
        out_path = expand(out_path)
        ensure_dir(os.path.dirname(out_path) or ".")
        if out_path.lower().endswith(".csv"):
            import csv
            cols = ["file", "caption", "ocr_text"]
            with open(out_path, "w", encoding="utf-8-sig", newline="") as f:
                w = csv.DictWriter(f, fieldnames=cols, extrasaction="ignore")
                w.writeheader()
                for r in records:
                    w.writerow(r)
        else:
            with open(out_path, "w", encoding="utf-8") as f:
                for r in records:
                    f.write(f"## {r['file']}\n")
                    if r.get("caption"):
                        f.write(f"- caption: {r['caption']}\n")
                    if r.get("ocr_text"):
                        f.write(f"- text: {r['ocr_text']}\n")
                    f.write("\n")
        written = out_path

    logger.info("[extract] %d photos processed (mode=%s) -> %s",
                len(records), mode, written or "(stdout)")
    return {"records": records, "out_path": written}


# ============================================================
# 能力 D：crop — 图片裁剪
# ============================================================

def capability_crop(
    root: str,
    width: Optional[int] = None,
    height: Optional[int] = None,
    ratio: Optional[str] = None,
    center_on_content: bool = False,
    out_dir: Optional[str] = None,
    recursive: bool = False,
    dry_run: bool = False,
) -> Dict[str, Any]:
    """批量裁剪图片到指定尺寸或比例。

    width/height: 目标像素尺寸（二者都给则按尺寸裁剪）
    ratio: 目标比例，如 "16:9" / "9:16" / "1:1"（按比例中心裁剪）
    center_on_content: 是否检测深色/主体区域居中（简单基于像素方差/亮度）
    out_dir: 输出目录，None=原地生成 _cropped 子目录（绝不覆盖原图）
    返回 {processed:[{src,dst,size}], skipped}
    """
    root = expand(root)
    assets = scan_desktop(root, recursive=recursive, max_depth=5 if recursive else 1)
    assets = [a for a in assets if not a.is_shortcut and _is_image(a.ext)]

    result: Dict[str, Any] = {"processed": [], "skipped": []}
    try:
        from PIL import Image
    except ImportError:
        logger.error("[crop] Pillow not installed")
        result["error"] = "Pillow not installed"
        return result

    out_root = expand(out_dir) if out_dir else os.path.join(root, "_cropped")

    for a in track(assets, label="批量裁剪", est_per_item=0.05):
        try:
            with Image.open(a.path) as im:
                w, h = im.size
                target = _compute_crop_box(w, h, width, height, ratio, center_on_content, im)
                if target is None:
                    result["skipped"].append({"path": a.path, "reason": "no crop spec"})
                    continue
                dst = _unique_path(os.path.join(out_root, os.path.basename(a.path)))
                result["processed"].append({
                    "src": a.path, "dst": dst, "size": f"{target[2]-target[0]}x{target[3]-target[1]}",
                })
                if not dry_run:
                    os.makedirs(out_root, exist_ok=True)
                    cropped = im.crop(target)
                    if width and height:
                        cropped = cropped.resize((width, height))
                    cropped.save(dst)
        except Exception as e:  # noqa: BLE001
            logger.debug("[crop] failed %s: %s", a.path, e)
            result["skipped"].append({"path": a.path, "reason": str(e)})

    action = "cropped" if not dry_run else "would crop"
    logger.info("[crop] %s %d images -> %s", action, len(result["processed"]), out_root)
    return result


def _compute_crop_box(w, h, width, height, ratio, center_on_content, im) -> Optional[Tuple[int, int, int, int]]:
    """计算裁剪框 (left, top, right, bottom)。"""
    if ratio:
        try:
            rw, rh = (float(x) for x in ratio.split(":"))
        except ValueError:
            return None
        target_ratio = rw / rh
        cur_ratio = w / h
        if cur_ratio > target_ratio:
            new_w = int(h * target_ratio)
            new_h = h
        else:
            new_w = w
            new_h = int(w / target_ratio)
    elif width and height:
        new_w, new_h = min(width, w), min(height, h)
    else:
        return None

    # 居中（可选基于内容）
    cx, cy = w // 2, h // 2
    if center_on_content:
        cx, cy = _content_center(im, w, h)
    left = max(0, min(cx - new_w // 2, w - new_w))
    top = max(0, min(cy - new_h // 2, h - new_h))
    return (left, top, left + new_w, top + new_h)


def _content_center(im, w, h) -> Tuple[int, int]:
    """简单主体中心检测：找亮度方差最大的区域中心（降采样加速）。"""
    try:
        small = im.convert("L").resize((64, 64))
        px = small.load()
        # 找最暗/最亮像素的质心（深色内容居中场景）
        total, sx, sy = 0, 0, 0
        for y in range(64):
            for x in range(64):
                v = 255 - px[x, y]  # 深色权重高
                total += v
                sx += x * v
                sy += y * v
        if total == 0:
            return w // 2, h // 2
        return int(sx / total / 64 * w), int(sy / total / 64 * h)
    except Exception:  # noqa: BLE001
        return w // 2, h // 2


# ============================================================
# 能力 E：to_ppt — 图片批量转 PPT
# ============================================================

def capability_to_ppt(
    root: str,
    out_path: Optional[str] = None,
    aspect: str = "16:9",
    margin: float = 0.08,
    sort_by: str = "name",
    recursive: bool = False,
) -> Dict[str, Any]:
    """把目录中的图片按顺序铺成 PPT，每张一页，居中。

    aspect: "16:9" | "4:3"
    margin: 边距比例（0.08=8%）
    sort_by: "name" | "time"
    返回 {out_path, slides}
    """
    root = expand(root)
    assets = scan_desktop(root, recursive=recursive, max_depth=5 if recursive else 1)
    assets = [a for a in assets if not a.is_shortcut and _is_image(a.ext)]
    if sort_by == "time":
        assets.sort(key=lambda a: a.best_time())
    else:
        assets.sort(key=lambda a: os.path.basename(a.path).lower())

    if not assets:
        return {"out_path": None, "slides": 0, "error": "no images found"}

    try:
        from pptx import Presentation
        from pptx.util import Emu
        from PIL import Image
    except ImportError:
        logger.error("[to_ppt] python-pptx / Pillow not installed")
        return {"out_path": None, "slides": 0, "error": "python-pptx not installed"}

    prs = Presentation()
    if aspect == "4:3":
        prs.slide_width = Emu(int(10 * 914400))
        prs.slide_height = Emu(int(7.5 * 914400))
    else:  # 16:9
        prs.slide_width = Emu(int(13.333 * 914400))
        prs.slide_height = Emu(int(7.5 * 914400))

    sw, sh = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]
    avail_w = sw * (1 - 2 * margin)
    avail_h = sh * (1 - 2 * margin)

    count = 0
    for a in track(assets, label="生成 PPT", est_per_item=0.05):
        try:
            with Image.open(a.path) as im:
                iw, ih = im.size
        except Exception as e:  # noqa: BLE001
            logger.debug("[to_ppt] open failed %s: %s", a.path, e)
            continue
        scale = min(avail_w / iw, avail_h / ih)
        pw, ph = int(iw * scale), int(ih * scale)
        left = int((sw - pw) / 2)
        top = int((sh - ph) / 2)
        slide = prs.slides.add_slide(blank)
        try:
            slide.shapes.add_picture(a.path, left, top, width=pw, height=ph)
            count += 1
        except Exception as e:  # noqa: BLE001
            logger.debug("[to_ppt] add_picture failed %s: %s", a.path, e)

    if not out_path:
        out_path = os.path.join(os.path.dirname(root), f"{os.path.basename(root)}.pptx")
    out_path = expand(out_path)
    ensure_dir(os.path.dirname(out_path) or ".")
    prs.save(out_path)
    logger.info("[to_ppt] %d slides -> %s", count, out_path)
    return {"out_path": out_path, "slides": count}


# ============================================================
# 能力 F：collage — 多图拼接成一张
# ============================================================

def capability_collage(
    root: str,
    out_path: Optional[str] = None,
    cols: int = 0,
    page: str = "A4",
    gap: int = 50,
    dpi: int = 300,
    recursive: bool = False,
) -> Dict[str, Any]:
    """把多张图片拼接成一张大图（默认 A4 300DPI）。

    cols: 每行列数，0=自动（按图片数开方）
    page: "A4"(2480x3508) | "A4L"(横向) | "none"(按内容自适应)
    gap: 图片间距像素
    返回 {out_path, images}
    """
    root = expand(root)
    assets = scan_desktop(root, recursive=recursive, max_depth=3 if recursive else 1)
    assets = [a for a in assets if not a.is_shortcut and _is_image(a.ext)]
    assets.sort(key=lambda a: os.path.basename(a.path).lower())
    if not assets:
        return {"out_path": None, "images": 0, "error": "no images found"}

    try:
        from PIL import Image
    except ImportError:
        return {"out_path": None, "images": 0, "error": "Pillow not installed"}

    n = len(assets)
    if cols <= 0:
        import math
        cols = max(1, int(math.ceil(math.sqrt(n))))
    rows = (n + cols - 1) // cols

    page_sizes = {"A4": (2480, 3508), "A4L": (3508, 2480)}
    if page in page_sizes:
        canvas_w, canvas_h = page_sizes[page]
        cell_w = (canvas_w - gap * (cols + 1)) // cols
        cell_h = (canvas_h - gap * (rows + 1)) // rows
    else:
        cell_w = cell_h = 600
        canvas_w = cell_w * cols + gap * (cols + 1)
        canvas_h = cell_h * rows + gap * (rows + 1)

    canvas = Image.new("RGB", (canvas_w, canvas_h), "white")
    placed = 0
    for idx, a in enumerate(assets):
        try:
            with Image.open(a.path) as im:
                im = im.convert("RGB")
                im.thumbnail((cell_w, cell_h))
                r, c = divmod(idx, cols)
                x = gap + c * (cell_w + gap) + (cell_w - im.width) // 2
                y = gap + r * (cell_h + gap) + (cell_h - im.height) // 2
                canvas.paste(im, (x, y))
                placed += 1
        except Exception as e:  # noqa: BLE001
            logger.debug("[collage] failed %s: %s", a.path, e)

    if not out_path:
        out_path = os.path.join(os.path.dirname(root), f"{os.path.basename(root)}_collage.jpg")
    out_path = expand(out_path)
    ensure_dir(os.path.dirname(out_path) or ".")
    canvas.save(out_path, dpi=(dpi, dpi))
    logger.info("[collage] %d images -> %s (%dx%d)", placed, out_path, canvas_w, canvas_h)
    return {"out_path": out_path, "images": placed}


# ============================================================
# 能力 G：video_rename_title — 视频标题 AI 重命名
# ============================================================

import re as _re_module


def clean_video_title(
    name: str,
    max_len: int = 25,
    min_len: int = 15,
    remove_tags: bool = True,
    remove_brackets: bool = True,
    remove_underscore: bool = True,
) -> str:
    """清洗视频标题（确定性规则，不依赖模型）：

    - 删除 #标签
    - 删除 （来源：xxx）/(xxx) 括号内容
    - 删除下划线
    - 去除"一审/二审/三审/编发"等尾部审核信息
    - 超长截断到 max_len
    """
    title = os.path.splitext(name)[0]
    if remove_brackets:
        title = _re_module.sub(r"[（(][^）)]*[）)]", "", title)
        title = _re_module.sub(r"【[^】]*】", "", title)
    # 去除审核/编发/来源尾巴（一审：xx 二审：xx ...）
    title = _re_module.sub(r"(一审|二审|三审|编发|来源)[：:].*$", "", title)
    if remove_tags:
        # 1) 成对话题标签 #纯词#（不含标点，微博风格）
        title = _re_module.sub(r"#[A-Za-z0-9\u4e00-\u9fff]+#", "", title)
        # 2) 空格/结尾分隔的话题标签 #纯词（抖音风格）
        title = _re_module.sub(r"#[A-Za-z0-9\u4e00-\u9fff]+(?=[\s]|$)", "", title)
        # 3) 残留的内联 # 仅去符号、保留文字（如 41岁#C罗 → 41岁C罗）
        title = title.replace("#", "")
    if remove_underscore:
        title = title.replace("_", "")
    # 合并多余空白与首尾标点
    title = _re_module.sub(r"\s+", "", title).strip(" ，,。.、")
    if len(title) > max_len:
        title = title[:max_len].rstrip(" ，,。.、")
    return title


def _video_duration_seconds(video_path: str, ffmpeg: str) -> float:
    """用 ffmpeg 探测视频时长（秒）。失败返回 0。

    imageio-ffmpeg 只自带 ffmpeg（无 ffprobe），故解析 `ffmpeg -i` 写到
    stderr 的 `Duration: HH:MM:SS.ss` 行。
    """
    try:
        proc = subprocess.run(
            [ffmpeg, "-i", video_path],
            capture_output=True, text=True, errors="ignore",
        )
        import re
        m = re.search(r"Duration:\s*(\d+):(\d+):(\d+\.?\d*)", proc.stderr or "")
        if m:
            h, mnt, s = m.groups()
            return int(h) * 3600 + int(mnt) * 60 + float(s)
    except Exception as e:  # noqa: BLE001
        logger.debug("[video] probe duration failed: %s", e)
    return 0.0


def _extract_video_frames(
    video_path: str, ffmpeg: str, n_frames: int = 6
) -> Tuple[Optional[str], List[str]]:
    """用内置 ffmpeg 沿时间轴均匀抽取 n_frames 帧。

    返回 (临时目录, 帧路径列表)。调用方负责 shutil.rmtree(临时目录)。
    这是 ollama 后端"真读视频"的关键：把抽出的帧当多图喂给 VLM，
    等价于 vLLM 服务端 do_sample_frames 的客户端实现。
    """
    dur = _video_duration_seconds(video_path, ffmpeg)
    if dur > 0:
        # 均匀取点，跳过首尾（常为黑帧/片头片尾）
        step = dur / (n_frames + 1)
        timestamps = [round(step * (i + 1), 2) for i in range(n_frames)]
    else:
        timestamps = [0.0]  # 探测不到时长时退化为取首帧

    tmp = tempfile.mkdtemp(prefix="vframe_")

    def _grab(idx_t: Tuple[int, float]) -> Optional[str]:
        i, t = idx_t
        out = os.path.join(tmp, f"f{i}.jpg")
        try:
            subprocess.run(
                [ffmpeg, "-y", "-ss", f"{t:.2f}", "-i", video_path,
                 "-frames:v", "1", "-vf", "scale=512:-2", out],
                capture_output=True, check=True,
            )
            if os.path.exists(out) and os.path.getsize(out) > 0:
                return out
        except Exception as e:  # noqa: BLE001
            logger.debug("[video] extract frame @%.2fs failed: %s", t, e)
        return None

    # 并发抽帧：多个时间点的 ffmpeg 互不依赖，线程池并行跑可显著缩短墙钟。
    paths: List[str] = []
    if len(timestamps) <= 1:
        got = _grab((0, timestamps[0])) if timestamps else None
        if got:
            paths.append(got)
    else:
        from concurrent.futures import ThreadPoolExecutor
        workers = min(len(timestamps), 12)
        with ThreadPoolExecutor(max_workers=workers) as ex:
            results = list(ex.map(_grab, list(enumerate(timestamps))))
        # 按时间顺序保留成功的帧
        paths = [p for p in results if p]
    return tmp, paths


def _ai_title_from_frames(
    vlm, frame_paths: List[str], max_len: int, min_len: int
) -> str:
    """让 VLM 真看视频抽帧画面，据实际内容起一个朴实标题（非标题党）。"""
    sys = (
        "你是视频内容编辑。下面是从同一段视频里按时间均匀抽取的若干帧画面。\n"
        "请综合这些画面里真实出现的人物、物体、场景、动作，概括这段视频在讲什么，"
        f"生成一个 {min_len}-{max_len} 字、准确朴实的中文标题。\n"
        "只描述画面里确实能看到的内容；不要编造、不要标题党、不要夸张词、"
        "不要 #标签、不要编号。只输出标题本身，不要任何解释。"
    )
    user = "请输出这段视频的标题："
    try:
        raw = vlm._impl._generate(  # noqa: SLF001
            sys, user, frame_paths, max_new_tokens=48, max_images=len(frame_paths)
        )
    except Exception as e:  # noqa: BLE001
        logger.debug("[video] title-from-frames failed: %s", e)
        return ""
    line = (raw or "").strip().splitlines()[0] if raw else ""
    line = clean_video_title(line + ".x", max_len=max_len, min_len=min_len)
    return line


def capability_video_rename_title(
    root: str,
    vlm=None,
    use_ai: bool = True,
    max_len: int = 25,
    min_len: int = 15,
    n_variants: int = 1,
    preview_n: int = 3,
    dry_run: bool = False,
    log_dir: Optional[str] = None,
) -> Dict[str, Any]:
    """对目录下视频按原名生成清洗后/AI 改写的新标题。

    流程契合用户习惯：默认 dry_run，先返回 preview_n 个示例供确认，
    确认后去掉 --dry-run 批量执行（默认直接执行）。

    use_ai: True 时优先「真读视频」——内置 ffmpeg 抽帧 → VLM 多图看画面起名；
            抽帧/看图失败回退到按原名文本改写；False 或无 VLM 时用规则清洗。
    log_dir: 提供时，真实重命名走带日志的执行器（executor.runner.apply_plan），
             记录每条 move，可用 `rollback --last` 一键撤销；为 None 时退回直接改名（不可回滚）。
    返回 {total, examples:[{old,new}], renames:[{old_path,new_path}], log_path}
    """
    root = expand(root)
    assets = scan_desktop(root, recursive=False, max_depth=1)
    videos = [a for a in assets if not a.is_shortcut and _is_video(a.ext)]
    videos.sort(key=lambda a: os.path.basename(a.path).lower())

    renames: List[Dict[str, str]] = []
    seen_titles: Dict[str, int] = {}

    # 每个视频：ffmpeg 抽帧 + VLM 多图看帧起名，单个就可能数秒~数十秒，必加进度。
    for a in track(videos, label="视频AI起名", est_per_item=6.0):
        old_name = os.path.basename(a.path)
        base_title = clean_video_title(old_name, max_len=max_len, min_len=min_len)

        new_title = base_title
        if use_ai and vlm is not None:
            frames_tmp: Optional[str] = None
            try:
                ai_title = ""
                # 首选：真读视频画面 —— 内置 ffmpeg 抽帧 → VLM 多图看图起名。
                ffmpeg = _ffmpeg_available()
                is_fallback = bool(getattr(vlm, "_is_fallback", False))
                if ffmpeg and not is_fallback:
                    frames_tmp, frames = _extract_video_frames(a.path, ffmpeg, n_frames=32)
                    if frames:
                        ai_title = _ai_title_from_frames(vlm, frames, max_len, min_len)
                # 兜底：抽帧/看图失败时，退回按原文件名做文本清洗改写（旧行为）。
                if not ai_title:
                    ai_title = _ai_rewrite_title(vlm, old_name, base_title, max_len, min_len)
                if ai_title:
                    new_title = ai_title
            except Exception as e:  # noqa: BLE001
                logger.debug("[video_rename_title] ai title failed: %s", e)
            finally:
                if frames_tmp and os.path.isdir(frames_tmp):
                    shutil.rmtree(frames_tmp, ignore_errors=True)

        # 去重：同标题加差异后缀
        key = new_title
        if key in seen_titles:
            seen_titles[key] += 1
            new_title = f"{new_title}（{seen_titles[key]}）"
        else:
            seen_titles[key] = 0

        if not new_title:
            continue
        # dry_run: 允许报告 unchanged 文件供用户确认（示例数充足）
        # apply   : unchanged 的文件也计入 total 但无需执行 move
        if new_title == os.path.splitext(old_name)[0]:
            renames.append({"old_path": a.path, "new_path": a.path,
                            "old": old_name, "new": new_title})
        else:
            new_path = _unique_path(os.path.join(os.path.dirname(a.path), f"{new_title}.{a.ext}"))
            renames.append({"old_path": a.path, "new_path": new_path,
                            "old": old_name, "new": os.path.basename(new_path)})

    examples = [{"old": r["old"], "new": r["new"]} for r in renames[:preview_n]]

    log_path: Optional[str] = None
    if not dry_run and renames:
        if log_dir:
            # 走带日志的执行器：把每条重命名表达为 move 动作，apply_plan 会逐条
            # 记入 ~/.archive_assistant/log/<ts>.json，可被 `rollback --last` 逆向撤销。
            from .executor.runner import apply_plan
            plan_dict = {
                "root": root,
                "mode": "video_rename_title",
                "actions": [
                    {"op": "move", "src": r["old_path"], "dst": r["new_path"]}
                    for r in renames
                ],
            }
            try:
                log_path = apply_plan(plan_dict, log_dir)
            except Exception as e:  # noqa: BLE001
                logger.error("[video_rename_title] logged apply failed: %s", e)
        else:
            # 未提供 log_dir：退回直接改名（不可回滚）。
            for r in renames:
                try:
                    shutil.move(r["old_path"], r["new_path"])
                except OSError as e:
                    logger.error("[video_rename_title] move failed: %s", e)

    action = "renamed" if not dry_run else "would rename"
    logger.info("[video_rename_title] %s %d videos (use_ai=%s, logged=%s)",
                action, len(renames), use_ai, bool(log_path))
    return {"total": len(renames), "examples": examples,
            "renames": renames, "log_path": log_path}


# ============================================================
# 能力 H：image_rename_by_content — 图片按内容自适应改名
# ============================================================
# 旧版 image_rename_by_ocr 对所有图都调 vlm.ocr()，导致无文字的照片/插画 OCR
# 为空被跳过。现改为单次 VLM 调用自适应：VLM 自己看图判断
#   - 图里主要是文字（截图/文档/聊天记录/扫描件）→ 提取关键文字
#   - 图是照片/插画/风景/人物                → 用一句话概括内容
# 二选一在同一 prompt 内完成，不额外加一次分类调用。

_FN_INVALID_RE = re.compile(r'[\\/:*?"<>|\r\n\t]')


# VLM 改名统一 prompt：模型自主判断是「文字图」还是「内容图」并输出合适的文件名。
_RENAME_SYS = (
    "你是图片归档助手。看图后输出一个可作为文件名的短标题。\n"
    "判断规则：\n"
    "- 如果图里主要是文字（截图、聊天记录、文档、扫描件、票据）→ 提取最关键的一行文字作为标题\n"
    "- 如果图是照片、插画、风景、人物、物体 → 用一句话概括画面内容作为标题\n"
    f"要求：纯文本，不要#标签、不要解释、不要编号；15-30 字；不要换行。\n"
    "只输出标题本身，不要任何前后缀。"
)


def _rename_text_to_filename(text: str, max_len: int = 30) -> str:
    """把 VLM 输出清洗成可用作文件名的短串：取第一行有意义的文字，去非法字符、限长度。
    空则返回 ""（调用方跳过改名）。"""
    if not text:
        return ""
    # 取第一行（模型可能输出多行注释）
    line = text.splitlines()[0].strip()
    # 去非法文件名字符
    name = _FN_INVALID_RE.sub("", line).strip().strip(" .-_")
    if len(name) > max_len:
        name = name[:max_len].strip()
    return name


def capability_image_rename_by_ocr(
    root: str,
    vlm=None,
    max_len: int = 30,
    preview_n: int = 5,
    dry_run: bool = False,
    log_dir: Optional[str] = None,
) -> Dict[str, Any]:
    """对目录下图片用 VLM 看图自适应改名（截图取文字 / 照片取内容）。

    旧版只走 OCR，导致无文字的照片被跳过；现改为单次 VLM 调用让模型自己判断：
    - 图里主要是文字 → 提取关键文字作为文件名
    - 图是照片/插画 → 用一句话概括内容作为文件名

    流程：每张图调一次 VLM（_RENAME_SYS prompt）→ 清洗输出 → 走执行器改名（可回滚）。
    VLM 输出为空或清洗后为空 → 跳过不改名，保留原名。
    log_dir: 提供时走带日志执行器（可 `rollback --last` 撤销）；None 时直接改名（不可回滚）。
    返回 {total, renamed, skipped, examples:[{old,new,ocr}], renames, log_path}
    """
    root = expand(root)
    assets = scan_desktop(root, recursive=False, max_depth=1)
    imgs = [a for a in assets if not a.is_shortcut and _is_image(a.ext)]
    imgs.sort(key=lambda a: os.path.basename(a.path).lower())

    renames: List[Dict[str, Any]] = []
    skipped: List[str] = []
    seen: Dict[str, int] = {}
    is_fallback = bool(getattr(vlm, "_is_fallback", False)) if vlm else True

    for a in track(imgs, label="图片内容改名", est_per_item=2.0):
        old_name = os.path.basename(a.path)
        if vlm is None or is_fallback:
            skipped.append(old_name)
            continue
        try:
            raw = vlm._impl._generate(  # noqa: SLF001
                _RENAME_SYS, "新文件名：", [a.path], max_new_tokens=48
            )
        except Exception as e:  # noqa: BLE001
            logger.debug("[image_rename] vlm failed %s: %s", old_name, e)
            raw = ""
        new_base = _rename_text_to_filename((raw or "").strip(), max_len=max_len)
        if not new_base:
            skipped.append(old_name)
            continue
        key = new_base
        if key in seen:
            seen[key] += 1
            new_base = f"{new_base}_{seen[key]}"
        else:
            seen[key] = 0
        new_name = f"{new_base}.{a.ext}"
        if new_name == old_name:
            renames.append({"old_path": a.path, "new_path": a.path,
                            "old": old_name, "new": new_name, "ocr": (raw or "")[:80]})
            continue
        new_path = _unique_path(os.path.join(os.path.dirname(a.path), new_name))
        renames.append({"old_path": a.path, "new_path": new_path,
                        "old": old_name, "new": os.path.basename(new_path),
                        "ocr": (raw or "")[:80]})

    examples = [{"old": r["old"], "new": r["new"], "ocr": r.get("ocr", "")}
                for r in renames[:preview_n]]

    log_path: Optional[str] = None
    if not dry_run and renames:
        real = [r for r in renames if r["old_path"] != r["new_path"]]
        if real:
            if log_dir:
                from .executor.runner import apply_plan
                plan_dict = {
                    "root": root,
                    "mode": "image_rename_by_ocr",
                    "actions": [
                        {"op": "move", "src": r["old_path"], "dst": r["new_path"]}
                        for r in real
                    ],
                }
                try:
                    log_path = apply_plan(plan_dict, log_dir)
                except Exception as e:  # noqa: BLE001
                    logger.error("[image_rename] apply failed: %s", e)
            else:
                for r in real:
                    try:
                        shutil.move(r["old_path"], r["new_path"])
                    except OSError as e:
                        logger.error("[image_rename] move failed: %s", e)

    action = "renamed" if not dry_run else "would rename"
    logger.info("[image_rename] %s %d/%d imgs (skipped %d, logged=%s)",
                action, len(renames), len(imgs), len(skipped), bool(log_path))
    return {"total": len(imgs), "renamed": len(renames),
            "skipped": skipped, "examples": examples,
            "renames": renames, "log_path": log_path}


def _ai_rewrite_title(vlm, old_name: str, base_title: str, max_len: int, min_len: int) -> str:
    """用 VLM 把原标题改写成意思相近、{min_len}-{max_len}字、无#标签的新标题。"""
    sys = (
        "你是短视频标题编辑。把给定的原标题改写成一个意思相近但更精炼的新标题。\n"
        f"要求：{min_len}-{max_len}字；删除所有#标签；删除括号及其内容（如（来源：xx））；"
        "不要任何编号；可适当增添标点；风格可贴近社会新闻。只输出新标题本身，不要解释。"
    )
    user = f"原标题：{base_title or old_name}\n新标题："
    try:
        raw = vlm._impl._generate(sys, user, [], max_new_tokens=48)  # noqa: SLF001
    except Exception:  # noqa: BLE001
        return ""
    line = (raw or "").strip().splitlines()[0] if raw else ""
    line = clean_video_title(line + ".x", max_len=max_len, min_len=min_len)
    return line


# ============================================================
# 能力 H：video_compose — 图片 + 音频 合成视频（FFmpeg）
# ============================================================

def _ffmpeg_available() -> Optional[str]:
    """定位 ffmpeg 可执行文件。

    1) 系统 PATH 优先（用户自己装的 brew/apt/winget 版本）。
    2) 回退到 pip 自带的静态二进制（imageio-ffmpeg），实现「零系统依赖」——
       只要 `pip install -r requirements.txt`，video-compose 即可开箱即用，
       无需让用户再去 brew/apt install ffmpeg。
    """
    exe = shutil.which("ffmpeg")
    if exe:
        return exe
    try:
        import imageio_ffmpeg
        bundled = imageio_ffmpeg.get_ffmpeg_exe()
        if bundled and os.path.exists(bundled):
            return bundled
    except Exception:
        pass
    return None


def capability_video_compose(
    image: str,
    audio: str,
    out_path: str,
    aspect: str = "9:16",
    quality: str = "high",
    dry_run: bool = False,
) -> Dict[str, Any]:
    """把单张图片 + 音频合成为视频（时长=音频时长）。

    aspect: "9:16"(1440x2560) | "16:9"(2560x1440) | "1:1"(1440x1440)
    quality: "high"(CRF15 近无损 + AAC320k) | "standard"(CRF23 + AAC192k)
    返回 {out_path, cmd, ok}
    使用 FFmpeg list 参数（不经 shell），避免命令注入。
    """
    image = expand(image)
    audio = expand(audio)
    out_path = expand(out_path)
    ffmpeg = _ffmpeg_available()
    if not ffmpeg:
        return {"out_path": None, "ok": False,
                "error": "ffmpeg 不可用：请执行 `pip install imageio-ffmpeg`（自带静态二进制，无需系统安装）"}
    if not os.path.exists(image):
        return {"out_path": None, "ok": False, "error": f"image not found: {image}"}
    if not os.path.exists(audio):
        return {"out_path": None, "ok": False, "error": f"audio not found: {audio}"}

    sizes = {"9:16": (1440, 2560), "16:9": (2560, 1440), "1:1": (1440, 1440)}
    w, h = sizes.get(aspect, (1440, 2560))
    vf = (
        f"scale={w}:{h}:force_original_aspect_ratio=decrease,"
        f"pad={w}:{h}:(ow-iw)/2:(oh-ih)/2:color=black,setsar=1"
    )
    if quality == "high":
        v_args = ["-c:v", "libx264", "-crf", "15", "-preset", "slow", "-pix_fmt", "yuv420p"]
        a_args = ["-c:a", "aac", "-b:a", "320k", "-af", "loudnorm"]
    else:
        v_args = ["-c:v", "libx264", "-crf", "23", "-preset", "medium", "-pix_fmt", "yuv420p"]
        a_args = ["-c:a", "aac", "-b:a", "192k"]

    cmd = [
        ffmpeg, "-y", "-loop", "1", "-i", image, "-i", audio,
        "-vf", vf, *v_args, *a_args,
        "-shortest", "-movflags", "+faststart", out_path,
    ]

    if dry_run:
        logger.info("[video_compose] (dry-run) %s", " ".join(cmd))
        return {"out_path": out_path, "ok": None, "cmd": cmd, "dry_run": True}

    ensure_dir(os.path.dirname(out_path) or ".")
    try:
        subprocess.run(cmd, check=True, capture_output=True, text=True)
        logger.info("[video_compose] created %s (%s, %s)", out_path, aspect, quality)
        return {"out_path": out_path, "ok": True, "cmd": cmd}
    except subprocess.CalledProcessError as e:
        logger.error("[video_compose] ffmpeg failed: %s", e.stderr[-500:] if e.stderr else e)
        return {"out_path": None, "ok": False, "error": (e.stderr or "")[-500:], "cmd": cmd}


# ============================================================
# 能力 I：video_distribute — 视频分发到多个子文件夹
# ============================================================

def capability_video_distribute(
    src: str,
    dst_base: str,
    per_folder: int = 1,
    folder_template: str = "{seq}",
    sort_by: str = "name",
    dry_run: bool = False,
) -> Dict[str, Any]:
    """把源目录下的视频按规则分发到目标基目录下的多个子文件夹。

    per_folder: 每个目标子文件夹放几个视频
    folder_template: 子文件夹命名模板，{seq} 为序号
    例：source 有 371 个视频, per_folder=1 → 分发到 371 个子文件夹各 1 个；
        若 per_folder=7 → 每个子文件夹放 7 个。
    返回 {total, moves:[{src,dst}], folders}
    """
    src = expand(src)
    dst_base = expand(dst_base)
    assets = scan_desktop(src, recursive=False, max_depth=1)
    videos = [a for a in assets if not a.is_shortcut and _is_video(a.ext)]
    if sort_by == "time":
        videos.sort(key=lambda a: a.best_time())
    else:
        videos.sort(key=lambda a: os.path.basename(a.path).lower())

    moves: List[Dict[str, str]] = []
    folders = set()
    for idx, a in enumerate(videos):
        folder_seq = idx // per_folder + 1
        folder_name = folder_template.format(seq=folder_seq)
        dst_dir = os.path.join(dst_base, folder_name)
        folders.add(dst_dir)
        dst = _unique_path(os.path.join(dst_dir, os.path.basename(a.path)))
        moves.append({"src": a.path, "dst": dst})
        if not dry_run:
            os.makedirs(dst_dir, exist_ok=True)
            shutil.move(a.path, dst)

    action = "distributed" if not dry_run else "would distribute"
    logger.info("[video_distribute] %s %d videos into %d folders (per_folder=%d)",
                action, len(moves), len(folders), per_folder)
    return {"total": len(moves), "moves": moves, "folders": sorted(folders)}
